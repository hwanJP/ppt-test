import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import io
import os
from datetime import datetime
from PIL import Image

def create_custom_slide(
    title="제목",
    subject="'주제'",
    paper_title="논문제목",
    eng_name="영문명",
    translation="- 내용 번역"
):
    """HK Kolmar 스타일 슬라이드 생성"""
    prs = Presentation()
    
    # 슬라이드 레이아웃 (빈 레이아웃 사용)
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 상단 바 (연회색)
    top_shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), prs.slide_width, Inches(0.6)  # type=1 (사각형)
    )
    fill = top_shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(242, 242, 242)
    top_shape.line.fill.background()  # 테두리 제거
    
    # 좌측 제목
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(3), Inches(0.5))
    title_box.text = title
    
    # 제목 텍스트 스타일 적용
    for paragraph in title_box.text_frame.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(0, 112, 192)
    
    # 우측 로고 자리 (로컬 logo.png 파일 사용)
    logo_path = "logo.png"
    if os.path.exists(logo_path):
        try:
            # 로고 이미지 추가
            slide.shapes.add_picture(
                logo_path, 
                prs.slide_width - Inches(2.2), 
                Inches(0.05), 
                height=Inches(0.4)  # 80% 크기
            )
        except Exception as e:
            st.warning(f"로고 삽입 중 오류: {str(e)}")
            # 오류 시 텍스트로 대체
            logo_text_box = slide.shapes.add_textbox(
                prs.slide_width - Inches(2.2), Inches(0.1), Inches(2), Inches(0.4)
            )
            logo_text_box.text = "HK Kolmar"
            
            # 로고 텍스트 스타일 적용
            for paragraph in logo_text_box.text_frame.paragraphs:
                paragraph.font.size = Pt(16)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(0, 112, 192)
    else:
        # 로고 파일이 없으면 텍스트로 표시
        logo_text_box = slide.shapes.add_textbox(
            prs.slide_width - Inches(2.2), Inches(0.1), Inches(2), Inches(0.4)
        )
        logo_text_box.text = "HK Kolmar"
        
        # 로고 텍스트 스타일 적용
        for paragraph in logo_text_box.text_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0, 112, 192)
    
    # 중앙 subject
    subject_box = slide.shapes.add_textbox(Inches(3), Inches(1.0), Inches(4), Inches(0.6))
    subject_box.text = subject
    
    # 주제 텍스트 스타일 적용
    for paragraph in subject_box.text_frame.paragraphs:
        paragraph.font.size = Pt(24)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(0, 112, 192)
    
    # 논문 제목
    paper_box = slide.shapes.add_textbox(Inches(3), Inches(1.6), Inches(6), Inches(0.6))
    paper_box.text = paper_title
    
    # 논문제목 텍스트 스타일 적용
    for paragraph in paper_box.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
    
    # SCIE 캡처 (우측 상단) - 예시 내용으로 박스 생성
    scie_box = slide.shapes.add_shape(1, Inches(7), Inches(1.0), Inches(2), Inches(1.2))
    scie_box.fill.solid()
    scie_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    scie_box.line.color.rgb = RGBColor(0, 0, 0)
    scie_box.line.width = Pt(1)
    
    # SCIE 박스 내용 작성 - 텍스트로 한번에 처리
    scie_box.text = "SCIE팀 논문\n\nSCIE 여부\n캘린더(등수)"
    
    # 텍스트 스타일 적용
    for paragraph in scie_box.text_frame.paragraphs:
        if paragraph.text == "SCIE팀 논문":
            paragraph.font.size = Pt(10)
            paragraph.font.bold = True
            paragraph.alignment = 1  # PP_ALIGN.CENTER
        else:
            paragraph.font.size = Pt(9)
            paragraph.alignment = 1  # PP_ALIGN.CENTER
    
    # 논문 인용 캡처 (중앙 큰 박스) - 예시 내용으로 채움
    citation_box = slide.shapes.add_shape(1, Inches(1.5), Inches(2.3), Inches(7), Inches(3))
    citation_box.fill.solid()
    citation_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    citation_box.line.color.rgb = RGBColor(0, 0, 0)
    citation_box.line.width = Pt(2)
    
    # 논문 인용 내용을 텍스트로 한번에 작성
    citation_text = """논문 인용 부분 하이라이트

• 연구 결과: 카멜리아 플라보노이드는 강력한 항산화 활성을 보임
• 효과: DPPH 라디칼 소거능 85.3% (500μg/mL)
• 메커니즘: ROS 생성 억제 및 항염 효과 확인
• 결론: 화장품 항산화 소재로 활용 가능성 높음

출처: Journal of Cosmetic Science, 2024, Vol.45, pp.123-135"""
    
    citation_box.text = citation_text
    
    # 텍스트 스타일 적용
    for i, paragraph in enumerate(citation_box.text_frame.paragraphs):
        if i == 0:  # 제목
            paragraph.font.size = Pt(16)
            paragraph.font.bold = True
            paragraph.alignment = 1  # PP_ALIGN.CENTER
            paragraph.font.color.rgb = RGBColor(0, 112, 192)
        elif "출처:" in paragraph.text:  # 출처
            paragraph.font.size = Pt(10)
            paragraph.font.italic = True
            paragraph.font.color.rgb = RGBColor(128, 128, 128)
            paragraph.alignment = 1  # PP_ALIGN.CENTER
        elif paragraph.text.strip():  # 내용
            paragraph.font.size = Pt(12)
    
    # 영문명
    eng_box = slide.shapes.add_textbox(Inches(0.3), Inches(5.5), Inches(12), Inches(0.5))
    eng_box.text = eng_name
    
    # 영문명 텍스트 스타일 적용 (가운데 정렬)
    for paragraph in eng_box.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = RGBColor(0, 112, 192)
        paragraph.alignment = 1  # PP_ALIGN.CENTER
    
    # 내용 번역
    trans_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(8), Inches(1))
    trans_box.text = translation
    
    # 내용 번역 텍스트 스타일 적용
    for paragraph in trans_box.text_frame.paragraphs:
        paragraph.font.size = Pt(14)
    
    # 하단 고정 문구
    footer_box = slide.shapes.add_textbox(
        Inches(0.3), prs.slide_height - Inches(0.5), Inches(9), Inches(0.5)
    )
    footer_box.text = "※ 본 자료는 단순 학술 자료로 검토되었습니다. 이외 국내 표시광고 실증 및 중국위생허가 등 규제에 대한 적절성은 별도의 추가 검토를 받으시기 바랍니다."
    
    # 하단 문구 텍스트 스타일 적용
    for paragraph in footer_box.text_frame.paragraphs:
        paragraph.font.size = Pt(10)
        paragraph.font.color.rgb = RGBColor(128, 128, 128)
    
    return prs

def create_content_slide(prs, slide_title, content_text, page_num):
    """내용 슬라이드 추가"""
    slide_layout = prs.slide_layouts[6]  # 빈 슬라이드
    slide = prs.slides.add_slide(slide_layout)
    
    # 상단 바 (연회색)
    top_shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), prs.slide_width, Inches(0.6)
    )
    fill = top_shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(242, 242, 242)
    top_shape.line.fill.background()
    
    # 우측 로고 (로컬 파일 사용)
    logo_path = "logo.png"
    if os.path.exists(logo_path):
        try:
            slide.shapes.add_picture(
                logo_path, 
                prs.slide_width - Inches(2.2), 
                Inches(0.05), 
                height=Inches(0.5)
            )
        except Exception as e:
            logo_text_box = slide.shapes.add_textbox(
                prs.slide_width - Inches(2.2), Inches(0.1), Inches(2), Inches(0.4)
            )
            logo_text_box.text = "HK Kolmar"
            
            # 로고 텍스트 스타일 적용
            for paragraph in logo_text_box.text_frame.paragraphs:
                paragraph.font.size = Pt(16)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(0, 112, 192)
    else:
        logo_text_box = slide.shapes.add_textbox(
            prs.slide_width - Inches(2.2), Inches(0.1), Inches(2), Inches(0.4)
        )
        logo_text_box.text = "HK Kolmar"
        
        # 로고 텍스트 스타일 적용
        for paragraph in logo_text_box.text_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0, 112, 192)
    
    # 슬라이드 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(11), Inches(1))
    title_box.text = slide_title
    
    # 제목 텍스트 스타일 적용
    for paragraph in title_box.text_frame.paragraphs:
        paragraph.font.size = Pt(28)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(0, 112, 192)
    
    # 내용 영역
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(11), Inches(5))
    content_box.text = content_text
    
    # 내용 텍스트 스타일 적용
    for paragraph in content_box.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = RGBColor(0, 0, 0)  # 검은색 명시
    
    # 페이지 번호
    page_box = slide.shapes.add_textbox(
        prs.slide_width - Inches(1), prs.slide_height - Inches(0.8), Inches(0.5), Inches(0.3)
    )
    page_box.text = str(page_num)
    
    # 페이지 번호 텍스트 스타일 적용
    for paragraph in page_box.text_frame.paragraphs:
        paragraph.font.size = Pt(14)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(0, 112, 192)
    
    return slide

def main():
    st.set_page_config(page_title="HK Kolmar PPT 생성기", page_icon="📄", layout="wide")
    
    st.title("📄 HK Kolmar 자동 로고 로드 PPT 생성기")
    st.markdown("같은 폴더의 logo.png 파일을 자동으로 로드하여 HK Kolmar 스타일 PPT를 생성합니다.")
    st.markdown("---")
    
    # 사이드바 설정
    with st.sidebar:
        st.header("📋 설정")
        
        # 로고 파일 상태 확인
        logo_path = "logo.png"
        if os.path.exists(logo_path):
            st.success("✅ logo.png 파일 발견")
            try:
                image = Image.open(logo_path)
                st.image(image, caption="현재 로고", width=150)
            except:
                st.warning("로고 파일을 읽을 수 없습니다")
        else:
            st.warning("⚠️ logo.png 파일이 없습니다")
            st.info("같은 폴더에 logo.png 파일을 넣어주세요")
        
        st.markdown("---")
        
        # 추가 슬라이드 설정
        include_content_slides = st.checkbox("내용 슬라이드 추가", value=False)
        num_content_slides = 0
        if include_content_slides:
            num_content_slides = st.slider("내용 슬라이드 수", 1, 5, 2)
    
    # 메인 입력 영역
    st.subheader("📝 기본 정보 입력")
    
    col1, col2 = st.columns([1, 1])
    
    with col1:
        title = st.text_input("제목 (좌상단)", value="안전성학술팀 보고서")
        subject = st.text_input("주제 (중앙 큰 제목)", value="'항산화 효능'", help="따옴표 포함")
        paper_title = st.text_input("논문제목", value="Flavonoids in Camellia japonica")
    
    with col2:
        eng_name = st.text_input("영문명", value="Camellia japonica flavonoids")
        translation = st.text_area(
            "내용 번역", 
            value="- 본 연구는 항산화 활성이 입증되었음을 보여줌",
            height=100
        )
    
    # 추가 내용 슬라이드 입력
    content_slides_data = []
    if include_content_slides:
        st.markdown("---")
        st.subheader("📄 추가 내용 슬라이드")
        
        for i in range(num_content_slides):
            st.markdown(f"### 내용 슬라이드 {i+1}")
            
            col1, col2 = st.columns([1, 2])
            
            with col1:
                slide_title = st.text_input(
                    f"슬라이드 제목 {i+1}",
                    value=f"연구 결과 {i+1}",
                    key=f"content_title_{i}"
                )
            
            with col2:
                slide_content = st.text_area(
                    f"슬라이드 내용 {i+1}",
                    value=f"• 주요 연구 결과 {i+1}\n• 데이터 분석 내용\n• 결론 및 제언",
                    height=100,
                    key=f"content_text_{i}"
                )
            
            content_slides_data.append({
                'title': slide_title,
                'content': slide_content
            })
    
    st.markdown("---")
    
    # PPT 생성
    st.subheader("🚀 PPT 생성")
    
    col1, col2, col3 = st.columns([1, 2, 1])
    
    with col2:
        if st.button("📄 HK Kolmar PPT 생성하기", type="primary", use_container_width=True):
            try:
                with st.spinner("HK Kolmar 스타일 PPT를 생성하는 중..."):
                    
                    # 메인 슬라이드 생성
                    prs = create_custom_slide(
                        title=title,
                        subject=subject,
                        paper_title=paper_title,
                        eng_name=eng_name,
                        translation=translation
                    )
                    
                    # 추가 내용 슬라이드들 생성
                    for i, slide_data in enumerate(content_slides_data):
                        create_content_slide(
                            prs,
                            slide_data['title'],
                            slide_data['content'],
                            i + 4  # 페이지 번호 (메인=3, 내용=4부터)
                        )
                    
                    # 파일로 저장
                    ppt_io = io.BytesIO()
                    prs.save(ppt_io)
                    ppt_io.seek(0)
                    
                    # 파일명 생성
                    safe_subject = subject.replace("'", "").replace(" ", "_")
                    filename = f"HK_Kolmar_{safe_subject}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
                    
                    st.success("✅ HK Kolmar PPT가 성공적으로 생성되었습니다!")
                    
                    # 다운로드 버튼
                    st.download_button(
                        label="📥 PPT 파일 다운로드",
                        data=ppt_io.getvalue(),
                        file_name=filename,
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True
                    )
                    
                    # 생성 정보
                    st.info(f"""
                    **생성된 PPT 정보:**
                    - 제목: {title}
                    - 주제: {subject}
                    - 논문: {paper_title}
                    - 로고: {'logo.png 사용' if os.path.exists('logo.png') else 'HK Kolmar 텍스트 사용'}
                    - 총 슬라이드: {1 + len(content_slides_data)}개
                    - 파일명: {filename}
                    """)
                    
            except Exception as e:
                st.error(f"❌ 오류가 발생했습니다: {str(e)}")
    
    # 사용 방법 안내
    with st.expander("📋 사용 방법 안내"):
        st.markdown("""
        ### 🔍 HK Kolmar 자동 로고 로드 시스템
        
        **1. 로고 파일 준비:**
        ```
        프로젝트 폴더/
        ├── streamlit_app.py (이 코드 파일)
        └── logo.png (HK Kolmar 로고)
        ```
        
        **2. 자동 로고 처리:**
        - `logo.png` 파일이 있으면 → 자동으로 이미지 삽입
        - `logo.png` 파일이 없으면 → "HK Kolmar" 텍스트 표시
        - 사이드바에서 로고 상태 실시간 확인
        
        **3. 예시 내용 자동 생성:**
        - **SCIE 박스**: "SCIE팀 논문", "SCIE 여부", "캘린더(등수)"
        - **논문 인용 박스**: 실제 연구 결과 예시
          - 카멜리아 플라보노이드 항산화 연구
          - DPPH 라디칼 소거능 데이터
          - ROS 생성 억제 메커니즘
          - 화장품 소재 활용 결론
          - 저널 출처 정보
        
        **4. 텍스트 서식:**
        - 모든 텍스트가 자연스럽게 줄바꿈 처리
        - 제목, 내용, 출처별 다른 스타일 적용
        - HK Kolmar 브랜드 색상 (파란색 #0070C0) 사용
        
        **5. 사용 팁:**
        - PNG 형식 로고 권장 (투명 배경)
        - 고해상도 이미지 사용 (300dpi 이상)
        - 주제에 따옴표 포함하여 입력 (예: '항산화 효능')
        - 내용 번역은 "- "로 시작하여 입력
        """)

if __name__ == "__main__":
    main()